import os
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
import traceback
from ai_logic import get_ai_response

def summarize_document(app, command, attached_file=None, **kwargs):
    if not attached_file or not attached_file.endswith('.docx'):
        return "Please attach a .docx file first, then use the command 'summarize this document'."
    try:
        app.speak_response(f"Reading {os.path.basename(attached_file)}...")
        doc = Document(attached_file)
        full_text = "\n".join([para.text for para in doc.paragraphs])
        if not full_text.strip(): return "The document appears to be empty."
        
        summary_prompt = f"Please provide a concise summary of the following document:\n\n{full_text[:15000]}"
        response_stream, _ = get_ai_response(app, [], summary_prompt)
        return "".join(list(response_stream))
    except Exception as e:
        app.queue_log(f"Failed to process document {attached_file}: {e}")
        return "I'm sorry, I ran into an error reading that document."

def create_document(app, doc_type, topic, **kwargs):
    app.speak_response(f"Okay, I'm researching and preparing your {doc_type} file about {topic}. This may take some time.")
    from skills.web_skill import perform_web_search
    
    try:
        # First, perform a web search to gather information
        search_results = perform_web_search(app, query=f"information about {topic}")
        if "error" in search_results.lower() or "couldn't find any web results" in search_results.lower():
             return "I couldn't find any information online to create your document."

        # Then, ask the AI to format this information
        formatting_prompt = (f"Based on the following information, please generate the content for a {doc_type} document "
                             f"about '{topic}'. Structure it appropriately with headings and paragraphs.\n\n"
                             f"--- INFORMATION ---\n{search_results}\n\n--- DOCUMENT CONTENT ---")

        app.queue_log(f"Generating content for {doc_type} on '{topic}' using AI...")
        response_stream, _ = get_ai_response(app, [], formatting_prompt)
        final_content = "".join(list(response_stream))

        if not final_content.strip():
            return "I wasn't able to generate content for your document. The AI did not produce any text."

        desktop_path = os.path.join(os.path.expanduser('~'), 'Desktop')
        # Ensure the directory exists
        os.makedirs(desktop_path, exist_ok=True) 
        
        safe_topic_name = "".join(c for c in topic if c.isalnum() or c in " ._").rstrip()
        file_path = "" # Initialize file_path

        if doc_type == "word":
            doc = Document()
            doc.add_heading(topic.title(), 0)
            # Split content by newlines and add as paragraphs to preserve basic formatting
            for paragraph in final_content.split('\n'):
                if paragraph.strip(): # Only add non-empty paragraphs
                    doc.add_paragraph(paragraph.strip())
            file_path = os.path.join(desktop_path, f"AURA_Word_{safe_topic_name}.docx")
            doc.save(file_path)
        elif doc_type == "powerpoint":
            prs = Presentation()
            slide_layout_title = prs.slide_layouts[0] # Title slide
            slide_layout_content = prs.slide_layouts[1] # Title and Content slide

            # Add title slide
            title_slide = prs.slides.add_slide(slide_layout_title)
            if title_slide.shapes.title:
                title_slide.shapes.title.text = topic.title()
            
            # Add content slide
            content_slide = prs.slides.add_slide(slide_layout_content)
            if content_slide.shapes.title:
                content_slide.shapes.title.text = "Summary and Key Points"
            if content_slide.shapes.placeholders:
                # Add text to the content placeholder, ensuring it fits
                text_frame = content_slide.shapes.placeholders[1].text_frame
                text_frame.clear() # Clear existing text
                # Add content line by line as paragraphs
                for paragraph in final_content.split('\n'):
                    if paragraph.strip():
                        p = text_frame.add_paragraph()
                        p.text = paragraph.strip()

            file_path = os.path.join(desktop_path, f"AURA_PPT_{safe_topic_name}.pptx")
            prs.save(file_path)
        else: # Handle 'excel' specifically or inform user it's not supported yet
            return f"Creating {doc_type} files is not yet supported."
        
        app.queue_log(f"Document saved to: {file_path}") # Log the actual save path
        return f"I've created your {doc_type} file about {topic} and saved it to your desktop at {file_path}."
    except Exception as e:
        app.queue_log(f"Failed to create document: {e}\n{traceback.format_exc()}", level="ERROR")
        return f"I ran into an error while creating the {doc_type} file. Details: {e}"

def register():
    return {
        'summarize_document': {
            'handler': summarize_document,
            'regex': r'summarize(?: this)? document',
            'params': [],
            'description': "Summarizes a .docx document that the user has attached."
        },
        'create_document': {
            'handler': create_document,
            'regex': r'create an? (word|powerpoint) (?:document|presentation|sheet|file) (?:on|about) (.+)', # Removed excel for now
            'params': ['doc_type', 'topic'],
            'description': "Creates a new Word or PowerPoint file on a topic by researching the web. Note: Excel file creation is not yet supported."
        },
    }