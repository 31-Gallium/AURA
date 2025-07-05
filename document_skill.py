# In skills/document_skill.py
import os
import re
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
import traceback
from ai_logic import run_ai_with_tools

def summarize_document(app, command, attached_file=None, **kwargs):
    """Reads an attached .docx file and uses AI to summarize it."""
    # This handler primarily relies on the file attached in the GUI.
    if not attached_file or not attached_file.endswith('.docx'):
        return "Please attach a .docx file first, then use the command 'summarize document'."

    try:
        app.speak_response(f"Reading {os.path.basename(attached_file)}...")
        doc = Document(attached_file)
        full_text = "\n".join([para.text for para in doc.paragraphs])
        if not full_text.strip(): return "The document appears to be empty."

        from ai_logic import get_ai_response
        summary_prompt = f"Please provide a concise summary of the following document:\n\n{full_text[:15000]}"
        summary = get_ai_response(app.answer_model, [], summary_prompt, app.queue_log)
        return summary
    except Exception as e:
        app.queue_log(f"Failed to process document {attached_file}: {e}")
        return "I'm sorry, I ran into an error reading that document."

def create_document(app, doc_type, topic, **kwargs):
    """Creates a Word, PowerPoint, or Excel file on a topic using AI."""
    # This handler now gets doc_type and topic directly from the command handler.
    user_prompt = (
        f"Create a structured and informative {doc_type} document about the topic: '{topic}'. "
        "Use your web search tool to gather accurate and up-to-date information before creating the content."
    )
    
    try:
        app.speak_response(f"Okay, I'm researching and preparing your {doc_type} file about {topic}. This may take some time.")
        final_content = run_ai_with_tools(app, user_prompt)

        if not final_content:
            return "I wasn't able to generate content for your document."

        desktop_path = os.path.join(os.path.expanduser('~'), 'Desktop')
        safe_topic_name = "".join(c for c in topic if c.isalnum() or c in " ._").rstrip()
        
        if doc_type == "word":
            doc = Document()
            doc.add_heading(topic.title(), 0)
            doc.add_paragraph(final_content)
            file_path = os.path.join(desktop_path, f"AURA_Word_{safe_topic_name}.docx")
            doc.save(file_path)
        elif doc_type == "powerpoint":
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = topic.title()
            text_slide = prs.slides.add_slide(prs.slide_layouts[1])
            text_slide.shapes.title.text = "Content Summary"
            text_slide.shapes.placeholders[1].text = final_content[:2000]
            file_path = os.path.join(desktop_path, f"AURA_PPT_{safe_topic_name}.pptx")
            prs.save(file_path)
        elif doc_type == "excel":
            wb = Workbook()
            ws = wb.active
            ws.title = topic[:30]
            for row in final_content.strip().split('\n'):
                ws.append([cell.strip() for cell in row.split(',')])
            file_path = os.path.join(desktop_path, f"AURA_Excel_{safe_topic_name}.xlsx")
            wb.save(file_path)
        
        return f"I've created your {doc_type} file about {topic} and saved it to your desktop."
    except Exception as e:
        app.queue_log(f"Failed to create document: {e}\n{traceback.format_exc()}")
        return f"I ran into an error while creating the {doc_type} file."

def register():
    """Registers document interaction commands with regex."""
    return {
        'summarize_document': {
            'handler': summarize_document,
            'regex': r'summarize(?: this)? document',
            'params': []
        },
        'create_document': {
            'handler': create_document,
            'regex': r'create an? (word|powerpoint|excel) (?:document|presentation|sheet|file) (?:on|about) (.+)',
            'params': ['doc_type', 'topic']
        },
    }